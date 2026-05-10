import argparse
import os
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    os.makedirs(os.path.dirname(args.output), exist_ok=True)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("min", [1, 2, 3])
    chart_data.add_series("max", [2, 3, 4])

    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(8),
        Inches(4.5),
        chart_data,
    )

    prs.save(args.output)


if __name__ == "__main__":
    main()
