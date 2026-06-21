"""Range chart range_name/legend wired into the spec layer (issue #14).

Non-valuation scenarios (allocation bands, spread ranges) must be able to
rename the legend and hide it from the spec/job layer, not hardcode "历史区间".
"""
import pandas as pd
from pptx import Presentation

from ablechart import render_chart, create_range_chart

_C = "{http://schemas.openxmlformats.org/drawingml/2006/chart}"


def _legend_texts(chart_el):
    out = []
    for s in chart_el.findall(f".//{_C}ser"):
        tx = s.find(f".//{_C}tx//{_C}v")
        if tx is not None and tx.text:
            out.append(tx.text)
    return out


def _df():
    return pd.DataFrame({"Asset": ["Equity", "Bond", "Gold"],
                         "Low": [2.5, 5.0, 1.0], "High": [15.0, 9.0, 6.0],
                         "Cur": [8.0, 7.0, 3.0]})


def test_spec_passes_range_name_override(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_chart(slide, {
        "chart": "range", "low": "Low", "high": "High", "current": "Cur",
        "range_name": "配置区间", "current_name": "当前权重", "format": "0.0%",
    }, _df())
    out = tmp_path / "r.pptx"
    prs.save(out)
    texts = _legend_texts(Presentation(str(out)).slides[0].shapes[0].chart._element)
    assert "配置区间" in texts
    assert "历史区间" not in texts


def test_spec_legend_none_hides_legend(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    render_chart(slide, {
        "chart": "range", "low": "Low", "high": "High",
        "range_name": "配置区间", "legend": "none",
    }, _df())
    out = tmp_path / "n.pptx"
    prs.save(out)
    chart = Presentation(str(out)).slides[0].shapes[0].chart
    assert chart.has_legend is False


def test_default_legend_name_still_works(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    create_range_chart(slide, _df(), "Asset", "Low", "High", current_col="Cur")
    out = tmp_path / "d.pptx"
    prs.save(out)
    texts = _legend_texts(Presentation(str(out)).slides[0].shapes[0].chart._element)
    assert "历史区间" in texts
